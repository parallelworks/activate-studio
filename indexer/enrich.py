#!/usr/bin/env python3
"""Enrich a GUFI index of the knowledge base with extracted text.

For every real-content file, extract plain text and store it in an fts5
`words` table inside the per-directory db.db of the GUFI index (joined to
entries by inode, per the GUFI master document pattern). DOCX/PDF/PPTX/XLSX
text is also written to an on-disk extract cache the server uses for
previews. Every directory db gets the `words` table, empty or not, so
MATCH queries never hit a missing table.
"""
import argparse
import os
import sqlite3
import subprocess
import sys
from pathlib import Path

EXCLUDE_DIRS = {
    '.git', 'node_modules', '.venv', 'venv', '__pycache__', 'dist', 'build',
    '.cache', '.pytest_cache', 'screenshots', '.ipynb_checkpoints',
}
TEXT_SUFFIXES = {
    '.md', '.txt', '.py', '.sh', '.yaml', '.yml', '.json', '.csv', '.tsv', '.xml',
    '.html', '.css', '.js', '.ts', '.tsx', '.svg', '.toml', '.cfg', '.ini', '.def', '.mjs', '.sql',
}
EXTRACT_SUFFIXES = {'.pdf', '.docx', '.pptx', '.xlsx'}
IMAGE_SUFFIXES = {'.png', '.jpg', '.jpeg', '.gif', '.webp'}
MAX_TEXT = 500_000
MAX_CAPTION_BYTES = 8 * 1024 * 1024
MIN_IMAGE_BYTES = 8 * 1024  # skip icons and tiny assets


def extract_docx(path: Path) -> str:
    from docx import Document  # type: ignore
    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append('\t'.join(c.text for c in row.cells))
    return '\n'.join(x for x in parts if x and x.strip())


def extract_pdf(path: Path) -> str:
    try:
        out = subprocess.run(['pdftotext', '-q', str(path), '-'],
                             capture_output=True, text=True, timeout=120)
        if out.returncode == 0 and out.stdout.strip():
            return out.stdout
    except FileNotFoundError:
        pass
    try:
        from pypdf import PdfReader  # type: ignore
        reader = PdfReader(str(path))
        return '\n'.join((page.extract_text() or '') for page in reader.pages)
    except Exception:
        return ''


def extract_pptx(path: Path) -> str:
    from pptx import Presentation  # type: ignore
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return '\n'.join(x for x in parts if x and x.strip())


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook  # type: ignore
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f'# sheet: {ws.title}')
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append('\t'.join(cells))
    return '\n'.join(parts)


def gateway_key() -> str:
    key = os.environ.get('PW_API_KEY', '')
    if key:
        return key
    try:
        import json
        import urllib.parse
        host = urllib.parse.urlparse(os.environ.get('PW_GATEWAY_URL', 'https://activate.parallel.works/api/openai/v1')).netloc
        creds = json.load(open(os.path.expanduser('~/.config/pw/credentials')))
        for ident in creds.get('identities', {}).values():
            if ident.get('server') == host:
                return ident.get('token') or ident.get('apikey') or ''
    except Exception:
        pass
    return ''


def caption_image(path: Path) -> str:
    """Describe an image with a vision model through the gateway's streaming
    Responses API. Opt-in via ADE_VISION_MODEL; failures degrade to OCR-only."""
    import base64
    import json
    import urllib.request
    model = os.environ.get('ADE_VISION_MODEL', '')
    key = gateway_key()
    if not model or not key or path.stat().st_size > MAX_CAPTION_BYTES:
        return ''
    mime = 'image/jpeg' if path.suffix.lower() in ('.jpg', '.jpeg') else f'image/{path.suffix.lower().lstrip(".")}'
    img = base64.b64encode(path.read_bytes()).decode()
    body = json.dumps({
        'model': model,
        'stream': True,
        'input': [{'role': 'user', 'content': [
            {'type': 'input_text', 'text': 'Describe this image in two to four factual sentences for a search index. Then transcribe any visible text labels.'},
            {'type': 'input_image', 'image_url': f'data:{mime};base64,{img}'}]}],
    }).encode()
    base = os.environ.get('PW_GATEWAY_URL', 'https://activate.parallel.works/api/openai/v1').rstrip('/')
    req = urllib.request.Request(f'{base}/responses', data=body,
                                 headers={'Authorization': f'Bearer {key}', 'Content-Type': 'application/json'})
    try:
        text = ''
        with urllib.request.urlopen(req, timeout=180) as resp:
            for raw in resp:
                line = raw.decode(errors='replace').strip()
                if not line.startswith('data:'):
                    continue
                data = line[5:].strip()
                if data == '[DONE]':
                    break
                try:
                    ev = json.loads(data)
                except Exception:
                    continue
                if ev.get('type') == 'response.output_text.delta':
                    text += ev.get('delta', '')
        return text.strip()
    except Exception as exc:
        print(f'  caption failed {path}: {str(exc)[:120]}', file=sys.stderr)
        return ''


def extract_image(path: Path) -> str:
    if path.stat().st_size < MIN_IMAGE_BYTES:
        return ''
    parts = []
    caption = caption_image(path)
    if caption:
        parts.append(f'Image description:\n{caption}')
    try:
        ocr = subprocess.run(['tesseract', str(path), 'stdout', '--psm', '3'],
                             capture_output=True, text=True, timeout=120)
        ocr_text = (ocr.stdout or '').strip()
        if len(ocr_text) > 20:
            parts.append(f'Text found in image (OCR):\n{ocr_text}')
    except Exception:
        pass
    return '\n\n'.join(parts)


def extract(path: Path) -> str | None:
    suffix = path.suffix.lower()
    try:
        if suffix in TEXT_SUFFIXES:
            return path.read_text(errors='replace')[:MAX_TEXT]
        if suffix == '.docx':
            return extract_docx(path)[:MAX_TEXT]
        if suffix == '.pdf':
            return extract_pdf(path)[:MAX_TEXT]
        if suffix == '.pptx':
            return extract_pptx(path)[:MAX_TEXT]
        if suffix == '.xlsx':
            return extract_xlsx(path)[:MAX_TEXT]
        if suffix in IMAGE_SUFFIXES:
            return extract_image(path)[:MAX_TEXT]
    except Exception as exc:
        print(f'  extract failed {path}: {exc}', file=sys.stderr)
    return None


def ensure_words_table(db: sqlite3.Connection) -> None:
    db.execute('DROP TABLE IF EXISTS words')
    db.execute("CREATE VIRTUAL TABLE words USING fts5(tinode UNINDEXED, fname UNINDEXED, wordf)")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument('--kb-root', required=True)
    ap.add_argument('--index', required=True, help='GUFI index tree top (mirrors kb-root)')
    ap.add_argument('--extract-cache', required=True)
    ap.add_argument('--subdir', default='', help='restrict to this KB-relative subtree (incremental indexing)')
    ap.add_argument('--no-recurse', action='store_true', help='process only the subdir itself, not its children')
    args = ap.parse_args()

    kb_root = Path(args.kb_root)
    index_root = Path(args.index)
    cache_root = Path(args.extract_cache)
    walk_root = index_root / args.subdir if args.subdir else index_root

    n_files = n_dirs = n_extracted = 0
    # Walk the index tree, not the source tree: every db.db the index holds
    # must carry the words table (empty or not) so MATCH queries never error.
    for dirpath, dirnames, _filenames in os.walk(walk_root):
        if args.no_recurse:
            dirnames[:] = []
        dirnames.sort()
        dbfile = Path(dirpath) / 'db.db'
        if not dbfile.exists():
            continue
        rel_dir = Path(dirpath).relative_to(index_root)
        src_dir = kb_root / rel_dir
        n_dirs += 1
        db = sqlite3.connect(str(dbfile))
        try:
            ensure_words_table(db)
            rows = []
            if src_dir.is_dir() and not any(part in EXCLUDE_DIRS for part in rel_dir.parts):
                for fpath in sorted(src_dir.iterdir()):
                    if not fpath.is_file():
                        continue
                    suffix = fpath.suffix.lower()
                    cacheable = suffix in EXTRACT_SUFFIXES or suffix in IMAGE_SUFFIXES
                    if suffix not in TEXT_SUFFIXES and not cacheable:
                        continue
                    # Expensive extractions (PDF, office, OCR, captions) are
                    # cached by mtime; a reindex pass reuses them untouched.
                    cache_file = cache_root / rel_dir / (fpath.name + '.txt')
                    if cacheable and cache_file.exists() and cache_file.stat().st_mtime >= fpath.stat().st_mtime:
                        text = cache_file.read_text(errors='replace')
                    else:
                        text = extract(fpath)
                        if cacheable:
                            # Cache even empty results so no-text images are
                            # not re-captioned on every pass.
                            cache_file.parent.mkdir(parents=True, exist_ok=True)
                            cache_file.write_text(text or '')
                            n_extracted += 1
                    if not text or not text.strip():
                        continue
                    rows.append((fpath.stat().st_ino, fpath.name, text))
                    n_files += 1
            db.executemany('INSERT INTO words (tinode, fname, wordf) VALUES (?, ?, ?)', rows)
            db.commit()
        finally:
            db.close()

    print(f'enriched {n_dirs} directories: {n_files} files indexed, {n_extracted} extract-cache entries')
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
